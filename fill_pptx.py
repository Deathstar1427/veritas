import sys
import os
from pptx import Presentation

def populate_presentation():
    input_file = "[EXT] Solution Challenge 2026 - Prototype PPT Template (1).pptx"
    output_file = "FairScan_Solution_Challenge_Presentation.pptx"
    
    if not os.path.exists(input_file):
        print(f"Error: {input_file} not found.")
        return

    prs = Presentation(input_file)

    slides_data = {
        "Team Details": "Team name: Death Dev\nTeam leader name: Anant Kumar\nProblem Statement: AI models in critical domains (hiring, loan, healthcare, education) are prone to bias. Detecting this requires deep technical expertise, making it a barrier for non-technical stakeholders to audit systems effectively.",
        "Brief about your solution": "FairScan (Veritas) is a full-stack, cloud-native fairness auditing web app. Users upload a dataset, and the system computes industry-standard bias metrics (DPD, EOD, DIR). It uses Google Gemini to generate plain-language explanations and mitigation guidance, and exports everything as a PDF audit report.",
        "Opportunities": "How different is it from any of the other existing ideas?\nExisting bias detection tools (like Fairlearn) are code-only libraries aimed at data scientists. FairScan provides a no-code visual interface for non-technical stakeholders.\n\nHow will it be able to solve the problem?\nBy automating complex statistical calculations and translating results into human-readable insights using Gemini.\n\nUSP of the proposed solution\nSeamless integration of rigorous statistical metrics with Google Gemini's generative AI to translate complex data into actionable, plain-English audit guidance.",
        "List of features offered by the solution": "• Domain-Specific Analysis (Hiring, Lending, Healthcare, Education)\n• Automated Bias Metrics (DPD, EOD, DIR)\n• Algorithmic Severity Classification\n• Generative AI Explanations via Google Gemini\n• Interactive Dashboard with Responsive Charts\n• Secure Access via Firebase Auth\n• Automated PDF Reporting",
        "Process flow diagram": "User Login -> Upload CSV -> FastAPI processes Fairlearn metrics -> Backend calls Gemini API -> React Dashboard Renders -> PDF Export\n\n(Please paste/insert your visual diagram over this text)",
        "Wireframes/Mock diagrams": "(Please insert screenshots of the FairScan Dashboard here)",
        "Architecture diagram": "React SPA on Firebase Hosting -> FastAPI on Google Cloud Run. Backend connects to Firebase Auth, Gemini API, and ReportLab\n\n(Please paste/insert your visual diagram over this text)",
        "Technologies to be used": "Frontend: React 19, Vite, Tailwind CSS, Tremor\nBackend: Python 3.11+, FastAPI, Uvicorn\nData/ML: Fairlearn, Pandas, Scikit-learn\nCloud: Firebase Hosting, Firebase Auth, Google Cloud Run\nAI: Google Gemini API (google-generativeai)",
        "Estimated implementation cost": "Development & Testing: $0 (Free Tiers)\nHosting & Auth (Firebase): $0 (Spark Plan)\nBackend (Cloud Run): Pay-per-use (Generous Free Tier)\nAI (Gemini API): Free tier access\nTotal MVP Cost: $0 to launch initially due to serverless architecture.",
        "Snapshots of the MVP": "(Please insert actual screenshots of FairScan: Login, Upload workspace, Results dashboard, PDF report here)",
        "Additional Details": "• Database Integration: Google Cloud Firestore for persisting audit histories.\n• Advanced Mitigation: Using Gemini to suggest data reweighing techniques.\n• Enterprise Features: Team workspaces and role-based access control (RBAC).",
        "Provide links": "GitHub Public Repository: [Insert Link]\nDemo Video Link (3 Minutes): [Insert Link]\nMVP Link: https://veritas-ai-01.web.app\nWorking Prototype Link: https://veritas-ai-01.web.app"
    }

    for slide in prs.slides:
        title_shape = slide.shapes.title if slide.shapes.title else None
        
        # Collect all text in the slide to identify it if title is missing or weird
        slide_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text += shape.text + " "
                
        # Find the best matching key
        matched_key = None
        for key in slides_data.keys():
            if key in slide_text:
                matched_key = key
                break
                
        if matched_key:
            # Try to find the body placeholder to put the text in
            body_shape = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    body_shape = shape
                    break
                    
            if body_shape and body_shape.has_text_frame:
                body_shape.text = slides_data[matched_key]
            else:
                # If no body placeholder, try to find the shape that currently holds the placeholder text
                for shape in slide.shapes:
                    if shape.has_text_frame and shape != title_shape:
                        if len(shape.text) > 10: # likely the content block
                            shape.text = slides_data[matched_key]
                            break

    prs.save(output_file)
    print(f"Successfully generated {output_file}")

if __name__ == '__main__':
    populate_presentation()
